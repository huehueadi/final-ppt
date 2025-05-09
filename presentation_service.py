from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from template_manager import TemplateManager
import requests
import json
import tempfile
import os
import logging
from PIL import Image
import io

logger = logging.getLogger(__name__)
template_manager = TemplateManager()
OLLAMA_ENDPOINT = "http://localhost:11434/api/generate"

def generate_text_content(topic, num_slides, custom_content=None):
    try:
        if custom_content:
            prompt = f"""Convert the following presentation content into a well-structured JSON format.
            The content provided by the user is about: '{topic}'
            
            USER CONTENT:
            {custom_content}
            
            Format EXACTLY as this JSON structure:
            {{
                "title": "Overall Presentation Title",
                "slides": [
                    {{
                        "title": "Slide 1 Title",
                        "points": [
                            "Point 1: Detailed explanation or context",
                            "Point 2: Detailed explanation or context",
                            "Point 3: Detailed explanation or context",
                            "Point 4: Additional context or related points",
                            "Point 5: Further insights or examples"
                        ]
                    }},
                    ...
                ]
            }}
            
            Requirements:
            - Use clear, professional language
            - Extract slide titles and bullet points from the user content
            - Organize the content logically
            - If the user hasn't provided enough structure, create appropriate slide titles and organize the content
            - Add bullet points where not explicitly provided by user
            - Avoid any markdown, code blocks, or extra formatting
            """
        else:
            prompt = f"""Generate a detailed JSON for a presentation about '{topic}' with {num_slides} slides.
            Each slide should have the following:
            - A detailed title
            - At least 5 concise and informative bullet points per slide (if applicable)
            - Provide some additional explanations or insights for each bullet point
            - Ensure the content is rich, professional, and informative
            Format EXACTLY as this JSON structure:
            {{
                "title": "Overall Presentation Title",
                "slides": [
                    {{
                        "title": "Slide 1 Title",
                        "points": [
                            "Point 1: Detailed explanation or context",
                            "Point 2: Detailed explanation or context",
                            "Point 3: Detailed explanation or context",
                            "Point 4: Additional context or related points",
                            "Point 5: Further insights or examples"
                        ]
                    }},
                    ...
                ]
            }}
            Requirements:
            - Use clear, professional language
            - Ensure each slide has a meaningful title
            - Create at least 5 detailed, informative bullet points per slide
            - Provide explanations, context, or examples where relevant
            - Avoid any markdown, code blocks, or extra formatting
            """
        
        payload = {
            "model": "gemma3:1b-it-qat",
            "prompt": prompt,
            "stream": False,
            "format": "json"
        }
        response = requests.post(OLLAMA_ENDPOINT, json=payload)
        if response.status_code != 200:
            logger.error(f"Ollama API error: {response.status_code} - {response.text}")
            raise Exception(f"Ollama API error: {response.status_code}")
        content = response.json()["response"]
        if "```json" in content:
            content = content.split("```json")[1].split("```")[0].strip()
        elif "```" in content:
            content = content.split("```")[1].split("```")[0].strip()
        presentation_data = json.loads(content)
        if not isinstance(presentation_data, dict) or 'title' not in presentation_data or 'slides' not in presentation_data:
            raise ValueError("Invalid JSON structure")
        for slide in presentation_data.get('slides', []):
            if 'title' not in slide or 'points' not in slide:
                raise ValueError("Invalid slide structure")
        return presentation_data
    except Exception as e:
        logger.error(f"Text generation error: {str(e)}")
        title = topic or "Presentation"
        return {
            "title": title,
            "slides": [
                {
                    "title": f"Introduction to {title}",
                    "points": [
                        "Overview of the topic with more context and background",
                        "Key points to discuss with additional details",
                        "Importance and relevance with examples or data"
                    ]
                },
                {
                    "title": "Main Concepts",
                    "points": [
                        "First main concept with detailed examples",
                        "Second main concept with further elaboration",
                        "Third main concept with supporting data or case studies"
                    ]
                },
                {
                    "title": "Conclusion",
                    "points": [
                        "Summary of key takeaways with insights",
                        "Future implications with potential applications",
                        "Call to action with a proposed next step or idea"
                    ]
                }
            ]
        }

def create_presentation(content_data, image_prompts=None, template="default"):
    try:
        template_config = template_manager.get_template(template) or template_manager.get_template('default')
        styles = template_config.get('styles', {})
        title_slide_styles = styles.get('title_slide', {})
        content_slide_styles = styles.get('content_slide', {})
        image_slide_styles = styles.get('image_slide', {})
        
        preview_image = template_config.get('preview_image', '')
        preview_image_path = os.path.join('static', preview_image)
        if preview_image and os.path.exists(preview_image_path):
            logger.info(f"Preview image found for template {template}: {preview_image_path}")
        else:
            logger.warning(f"Preview image not found for template {template}: {preview_image_path}")
        
        preview_data = {
            "title": content_data.get("title", "Presentation"),
            "template": template,
            "styles": {
                "title_slide": {
                    "background": title_slide_styles.get('background', {'type': 'solid', 'color': {'r': 240, 'g': 240, 'b': 240}}),
                    "background_image": title_slide_styles.get('background_image', ''),
                    "title_font": title_slide_styles.get('title_font', {'name': 'Calibri', 'size': 44, 'color': {'r': 0, 'g': 0, 'b': 0}, 'bold': True, 'alignment': 'center'}),
                    "image_position": title_slide_styles.get('image_position', {'left': 2.5, 'top': 4.0, 'width': 5.0, 'height': 2.5})
                },
                "content_slide": {
                    "background": content_slide_styles.get('background', {'type': 'solid', 'color': {'r': 255, 'g': 255, 'b': 255}}),
                    "background_image": content_slide_styles.get('background_image', ''),
                    "title_font": content_slide_styles.get('title_font', {'name': 'Calibri', 'size': 32, 'color': {'r': 0, 'g': 0, 'b': 0}, 'bold': True, 'alignment': 'left'}),
                    "body_font": content_slide_styles.get('body_font', {'name': 'Calibri', 'size': 18, 'color': {'r': 50, 'g': 50, 'b': 50}, 'alignment': 'left'}),
                    "image_position": content_slide_styles.get('image_position', {'left': 6.0, 'top': 1.5, 'width': 3.5, 'height': 4.5})
                },
                "image_slide": {
                    "fill_color": image_slide_styles.get('fill_color', {'r': 245, 'g': 245, 'b': 245}),
                    "border_color": image_slide_styles.get('border_color', {'r': 200, 'g': 200, 'b': 200}),
                    "border_width": image_slide_styles.get('border_width', 1.5),
                    "border_style": image_slide_styles.get('border_style', 'dashed')
                }
            },
            "slides": []
        }
        
        prs = Presentation()
        SLIDE_WIDTH = Inches(10)
        SLIDE_HEIGHT = Inches(7.5)
        SUPPORTED_FORMATS = {'BMP', 'GIF', 'JPEG', 'PNG', 'TIFF', 'WMF'}
        
        def validate_image_format(image_path):
            try:
                with Image.open(image_path) as img:
                    format = img.format.upper()
                    if format not in SUPPORTED_FORMATS:
                        logger.warning(f"Unsupported image format at {image_path}: got {format}, expected one of {SUPPORTED_FORMATS}")
                        return False
                    logger.debug(f"Validated image format at {image_path}: {format}")
                    return True
            except Exception as e:
                logger.error(f"Failed to validate image format at {image_path}: {str(e)}")
                return False
        
        def adjust_font_size(title_text, base_size):
            if len(title_text) > 40:
                new_size = max(base_size - 8, 20)
                logger.debug(f"Reducing font size for title '{title_text[:20]}...': {base_size}pt to {new_size}pt")
                return new_size
            return base_size
        
        blank_slide_layout = prs.slide_layouts[6]
        title_slide = prs.slides.add_slide(blank_slide_layout)
        
        background_settings = title_slide_styles.get('background', {})
        bg_image = title_slide_styles.get('background_image', '')
        background = title_slide.background
        fill = background.fill
        if bg_image:
            bg_image_path = os.path.abspath(os.path.join('static', bg_image))
            logger.debug(f"Checking background image for title slide: {bg_image_path}")
            if os.path.exists(bg_image_path) and validate_image_format(bg_image_path):
                try:
                    logger.info(f"Applying background image for title slide: {bg_image_path}")
                    picture = title_slide.shapes.add_picture(
                        bg_image_path,
                        left=0,
                        top=0,
                        width=SLIDE_WIDTH,
                        height=SLIDE_HEIGHT
                    )
                    logger.debug(f"Picture added to title slide: width={picture.width.inches:.2f}in, height={picture.height.inches:.2f}in, left={picture.left.inches:.2f}in, top={picture.top.inches:.2f}in")
                    title_slide.shapes._spTree.remove(picture._element)
                    title_slide.shapes._spTree.insert(2, picture._element)
                except Exception as e:
                    logger.error(f"Failed to apply background image for title slide: {bg_image_path}, error: {str(e)}")
                    fill.solid()
                    bg_color = background_settings.get('color', {'r': 240, 'g': 240, 'b': 240}) if background_settings.get('type') == 'solid' else background_settings.get('gradient_start', {'r': 240, 'g': 240, 'b': 240})
                    fill.fore_color.rgb = RGBColor(bg_color['r'], bg_color['g'], bg_color['b'])
                    logger.info(f"Fallback to solid color for title slide: rgb({bg_color['r']}, {bg_color['g']}, {bg_color['b']})")
            else:
                logger.error(f"Background image not found or invalid format for title slide: {bg_image_path}")
                fill.solid()
                bg_color = background_settings.get('color', {'r': 240, 'g': 240, 'b': 240}) if background_settings.get('type') == 'solid' else background_settings.get('gradient_start', {'r': 240, 'g': 240, 'b': 240})
                fill.fore_color.rgb = RGBColor(bg_color['r'], bg_color['g'], bg_color['b'])
                logger.info(f"Fallback to solid color for title slide: rgb({bg_color['r']}, {bg_color['g']}, {bg_color['b']})")
        else:
            logger.debug(f"No background image specified for title slide, using {background_settings.get('type', 'solid')} background")
            fill.solid()
            bg_color = background_settings.get('color', {'r': 240, 'g': 240, 'b': 240}) if background_settings.get('type') == 'solid' else background_settings.get('gradient_start', {'r': 240, 'g': 240, 'b': 240})
            fill.fore_color.rgb = RGBColor(bg_color['r'], bg_color['g'], bg_color['b'])
            logger.info(f"Applied solid color for title slide: rgb({bg_color['r']}, {bg_color['g']}, {bg_color['b']})")
        
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9.0)
        height = Inches(2.0)
        title_box = title_slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_text = content_data.get("title", "Presentation")
        title_frame.text = title_text
        logger.debug(f"Title slide heading: '{title_text}', length: {len(title_text)}")
        title_para = title_frame.paragraphs[0]
        title_font_settings = title_slide_styles.get('title_font', {})
        title_para.font.name = title_font_settings.get('name', 'Calibri')
        base_font_size = title_font_settings.get('size', 44)
        title_para.font.size = Pt(adjust_font_size(title_text, base_font_size))
        title_color = title_font_settings.get('color', {'r': 0, 'g': 0, 'b': 0})
        title_para.font.color.rgb = RGBColor(title_color['r'], title_color['g'], title_color['b'])
        title_para.font.bold = title_font_settings.get('bold', True)
        title_para.alignment = PP_ALIGN.CENTER
        
        title_image_style = {}
        if image_prompts and "title" in image_prompts:
            image_position = title_slide_styles.get('image_position', {'left': 2.5, 'top': 4.0, 'width': 5.0, 'height': 2.5})
            img_left = Inches(image_position.get('left', 2.5))
            img_top = Inches(image_position.get('top', 4.0))
            img_width = Inches(image_position.get('width', 5.0))
            img_height = Inches(image_position.get('height', 2.5))
            img_placeholder = title_slide.shapes.add_shape(1, img_left, img_top, img_width, img_height)
            img_placeholder.fill.solid()
            fill_color = image_slide_styles.get('fill_color', {'r': 245, 'g': 245, 'b': 245})
            img_placeholder.fill.fore_color.rgb = RGBColor(fill_color['r'], fill_color['g'], fill_color['b'])
            border_color = image_slide_styles.get('border_color', {'r': 200, 'g': 200, 'b': 200})
            img_placeholder.line.color.rgb = RGBColor(border_color['r'], border_color['g'], border_color['b'])
            img_placeholder.line.width = Pt(image_slide_styles.get('border_width', 1.5))
            img_placeholder.line.dash_style = 2 if image_slide_styles.get('border_style', 'dashed') == 'dashed' else 1
            text_frame = img_placeholder.text_frame
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            icon_p = text_frame.add_paragraph()
            icon_p.text = "🖼️"
            icon_p.alignment = PP_ALIGN.CENTER
            icon_p.font.size = Pt(48)
            icon_p.space_after = Pt(10)
            prompt_p = text_frame.add_paragraph()
            prompt_p.text = image_prompts['title']
            prompt_p.alignment = PP_ALIGN.CENTER
            prompt_p.font.italic = True
            prompt_p.font.size = Pt(14)
            prompt_p.font.color.rgb = RGBColor(100, 100, 100)
            title_image_style = {
                "left": image_position.get('left', 2.5),
                "top": image_position.get('top', 4.0),
                "width": image_position.get('width', 5.0),
                "height": image_position.get('height', 2.5),
                "fill_color": fill_color,
                "border_color": border_color,
                "border_width": image_slide_styles.get('border_width', 1.5),
                "border_style": image_slide_styles.get('border_style', 'dashed')
            }
        
        preview_data["slides"].append({
            "type": "title",
            "title": content_data.get("title", "Presentation"),
            "has_image": "title" in image_prompts if image_prompts else False,
            "image_prompt": image_prompts.get("title") if image_prompts else None,
            "image_style": title_image_style
        })
        
        for i, slide_data in enumerate(content_data.get("slides", [])):
            slide_index = str(i)
            content_slide = prs.slides.add_slide(blank_slide_layout)
            background_settings = content_slide_styles.get('background', {})
            bg_image = content_slide_styles.get('background_image', '')
            background = content_slide.background
            fill = background.fill
            if bg_image:
                bg_image_path = os.path.abspath(os.path.join('static', bg_image))
                logger.debug(f"Checking background image for content slide {i+1}: {bg_image_path}")
                if os.path.exists(bg_image_path) and validate_image_format(bg_image_path):
                    try:
                        logger.info(f"Applying background image for content slide {i+1}: {bg_image_path}")
                        picture = content_slide.shapes.add_picture(
                            bg_image_path,
                            left=0,
                            top=0,
                            width=SLIDE_WIDTH,
                            height=SLIDE_HEIGHT
                        )
                        logger.debug(f"Picture added to content slide {i+1}: width={picture.width.inches:.2f}in, height={picture.height.inches:.2f}in, left={picture.left.inches:.2f}in, top={picture.top.inches:.2f}in")
                        content_slide.shapes._spTree.remove(picture._element)
                        content_slide.shapes._spTree.insert(2, picture._element)
                    except Exception as e:
                        logger.error(f"Failed to apply background image for content slide {i+1}: {bg_image_path}, error: {str(e)}")
                        fill.solid()
                        bg_color = background_settings.get('color', {'r': 255, 'g': 255, 'b': 255}) if background_settings.get('type') == 'solid' else background_settings.get('gradient_start', {'r': 255, 'g': 255, 'b': 255})
                        fill.fore_color.rgb = RGBColor(bg_color['r'], bg_color['g'], bg_color['b'])
                        logger.info(f"Fallback to solid color for content slide {i+1}: rgb({bg_color['r']}, {bg_color['g']}, {bg_color['b']})")
                else:
                    logger.error(f"Background image not found or invalid format for content slide {i+1}: {bg_image_path}")
                    fill.solid()
                    bg_color = background_settings.get('color', {'r': 255, 'g': 255, 'b': 255}) if background_settings.get('type') == 'solid' else background_settings.get('gradient_start', {'r': 255, 'g': 255, 'b': 255})
                    fill.fore_color.rgb = RGBColor(bg_color['r'], bg_color['g'], bg_color['b'])
                    logger.info(f"Fallback to solid color for content slide {i+1}: rgb({bg_color['r']}, {bg_color['g']}, {bg_color['b']})")
            else:
                logger.debug(f"No background image specified for content slide {i+1}, using {background_settings.get('type', 'solid')} background")
                fill.solid()
                bg_color = background_settings.get('color', {'r': 255, 'g': 255, 'b': 255}) if background_settings.get('type') == 'solid' else background_settings.get('gradient_start', {'r': 255, 'g': 255, 'b': 255})
                fill.fore_color.rgb = RGBColor(bg_color['r'], bg_color['g'], bg_color['b'])
                logger.info(f"Applied solid color for content slide {i+1}: rgb({bg_color['r']}, {bg_color['g']}, {bg_color['b']})")
            
            title_left = Inches(0.5)
            title_top = Inches(0.5)
            title_width = Inches(9.0)
            title_height = Inches(1.2)
            title_box = content_slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_text = slide_data.get("title", f"Slide {i+1}")
            title_frame.text = title_text
            logger.debug(f"Content slide {i+1} heading: '{title_text}', length: {len(title_text)}")
            title_para = title_frame.paragraphs[0]
            title_font = content_slide_styles.get('title_font', {})
            title_para.font.name = title_font.get('name', 'Calibri')
            base_font_size = title_font.get('size', 32)
            title_para.font.size = Pt(adjust_font_size(title_text, base_font_size))
            title_color = title_font.get('color', {'r': 0, 'g': 0, 'b': 0})
            title_para.font.color.rgb = RGBColor(title_color['r'], title_color['g'], title_color['b'])
            title_para.font.bold = title_font.get('bold', True)
            title_para.alignment = {
                'center': PP_ALIGN.CENTER,
                'left': PP_ALIGN.LEFT,
                'right': PP_ALIGN.RIGHT
            }.get(title_font.get('alignment', 'left'), PP_ALIGN.LEFT)
            
            points_styling = []
            if slide_data.get("points", []):
                content_left = Inches(0.5)
                content_top = Inches(2.0)
                content_width = Inches(5.0)
                content_height = Inches(5.0)
                content_box = content_slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
                text_frame = content_box.text_frame
                text_frame.word_wrap = True
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                body_font = content_slide_styles.get('body_font', {})
                for point in slide_data.get("points", []):
                    if text_frame.paragraphs and text_frame.paragraphs[0].text == "":
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = "• " + point
                    p.font.name = body_font.get('name', 'Calibri')
                    p.font.size = Pt(body_font.get('size', 18))
                    body_color = body_font.get('color', {'r': 50, 'g': 50, 'b': 50})
                    p.font.color.rgb = RGBColor(body_color['r'], body_color['g'], body_color['b'])
                    p.space_before = Pt(6)
                    p.space_after = Pt(6)
                    p.alignment = {
                        'center': PP_ALIGN.CENTER,
                        'left': PP_ALIGN.LEFT,
                        'right': PP_ALIGN.RIGHT
                    }.get(body_font.get('alignment', 'left'), PP_ALIGN.LEFT)
                    points_styling.append({
                        "text": point,
                        "level": 0,
                        "font_name": body_font.get('name', 'Calibri'),
                        "font_size": body_font.get('size', 18),
                        "color": body_color,
                        "alignment": body_font.get('alignment', 'left'),
                        "space_before": 6,
                        "space_after": 6
                    })
            
            content_image_style = {}
            if image_prompts and slide_index in image_prompts:
                image_position = content_slide_styles.get('image_position', {'left': 6.0, 'top': 2.0, 'width': 3.5, 'height': 4.0})
                img_left = Inches(image_position.get('left', 6.0))
                img_top = Inches(image_position.get('top', 2.0))
                img_width = Inches(image_position.get('width', 3.5))
                img_height = Inches(image_position.get('height', 4.0))
                img_placeholder = content_slide.shapes.add_shape(1, img_left, img_top, img_width, img_height)
                img_placeholder.fill.solid()
                fill_color = image_slide_styles.get('fill_color', {'r': 245, 'g': 245, 'b': 245})
                img_placeholder.fill.fore_color.rgb = RGBColor(fill_color['r'], fill_color['g'], fill_color['b'])
                border_color = image_slide_styles.get('border_color', {'r': 200, 'g': 200, 'b': 200})
                img_placeholder.line.color.rgb = RGBColor(border_color['r'], border_color['g'], border_color['b'])
                img_placeholder.line.width = Pt(image_slide_styles.get('border_width', 1.5))
                img_placeholder.line.dash_style = 2 if image_slide_styles.get('border_style', 'dashed') == 'dashed' else 1
                text_frame = img_placeholder.text_frame
                text_frame.word_wrap = True
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                icon_p = text_frame.add_paragraph()
                icon_p.text = "🖼️"
                icon_p.alignment = PP_ALIGN.CENTER
                icon_p.font.size = Pt(48)
                icon_p.space_after = Pt(10)
                prompt_p = text_frame.add_paragraph()
                prompt_p.text = image_prompts[slide_index]
                prompt_p.alignment = PP_ALIGN.CENTER
                prompt_p.font.italic = True
                prompt_p.font.size = Pt(14)
                prompt_p.font.color.rgb = RGBColor(100, 100, 100)
                content_image_style = {
                    "left": image_position.get('left', 6.0),
                    "top": image_position.get('top', 2.0),
                    "width": image_position.get('width', 3.5),
                    "height": image_position.get('height', 4.0),
                    "fill_color": fill_color,
                    "border_color": border_color,
                    "border_width": image_slide_styles.get('border_width', 1.5),
                    "border_style": image_slide_styles.get('border_style', 'dashed')
                }
            
            preview_data["slides"].append({
                "type": "content",
                "title": slide_data.get("title", f"Slide {i+1}"),
                "points": slide_data.get("points", []),
                "points_styling": points_styling,
                "has_image": slide_index in image_prompts if image_prompts else False,
                "image_prompt": image_prompts.get(slide_index) if image_prompts else None,
                "image_style": content_image_style
            })
        
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        prs.save(temp_file.name)
        temp_file.close()
        return temp_file.name, preview_data
    
    except Exception as e:
        logger.error(f"PowerPoint creation error: {str(e)}")
        raise Exception(f"Failed to create PowerPoint: {str(e)}")

def generate_slide_previews(pptx_path):
    try:
        prs = Presentation(pptx_path)
        slide_previews = []
        
        for slide_idx, slide in enumerate(prs.slides):
            # Initialize default values
            title = slide.shapes.title.text if slide.shapes.title else f"Slide {slide_idx + 1}"
            content = ""
            
            # Create a blank image
            img = Image.new('RGB', (960, 540), color=(255, 255, 255))
            from PIL import ImageDraw, ImageFont
            draw = ImageDraw.Draw(img)
            try:
                font = ImageFont.load_default()
                # Collect text from shapes that have a text frame, excluding the title
                text_shapes = [
                    shape.text for shape in slide.shapes
                    if shape.has_text_frame and shape.text and shape != slide.shapes.title
                ]
                content = "\n".join(text_shapes) if text_shapes else "No content available"
                
                # Render title and content on the image
                draw.text((50, 50), title, fill=(0, 0, 0), font=font)
                y_text = 100
                for line in content.split('\n'):
                    draw.text((50, y_text), line, fill=(50, 50, 50), font=font)
                    y_text += 30
            except Exception as e:
                logger.error(f"Error rendering text for slide {slide_idx + 1}: {str(e)}")
                content = "Error rendering slide content"
                draw.text((50, 100), content, fill=(50, 50, 50), font=font)
            
            # Save image to BytesIO
            img_byte_arr = io.BytesIO()
            img.save(img_byte_arr, format='PNG')
            img_byte_arr.seek(0)
            slide_previews.append({
                'image': img_byte_arr,
                'title': title,
                'content': content
            })
        
        return slide_previews
    except Exception as e:
        logger.error(f"Error generating slide previews for {pptx_path}: {str(e)}")
        raise Exception(f"Failed to generate slide previews: {str(e)}")